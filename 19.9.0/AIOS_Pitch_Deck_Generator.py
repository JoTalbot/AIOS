import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = pptx.Presentation()

    # Слайд 1: Титульный
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AIOS v19.0.0"
    subtitle.text = "Self-Evolving Enterprise AI\nThe Skynet Epoch"

    # Слайд 2: Проблема
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Проблема: Статичный ИИ"
    content = slide.placeholders[1]
    content.text = "• Современные фреймворки ИИ ломаются при сбоях API.\n• ИИ-агенты забывают контекст после сессии (отсутствие долгосрочной памяти).\n• Кодовая база требует постоянной поддержки людьми.\n• Интеграции и платформы устаревают каждый день."

    # Слайд 3: Решение - AIOS Swarm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Решение: AIOS v19.0.0"
    content = slide.placeholders[1]
    content.text = "• Мультиагентный рой: Архитектор, Страж, Мета-Кодер.\n• Deep RAG Memory (ChromaDB) — агенты помнят весь прошлый опыт.\n• Model Context Protocol (MCP) — глаза (Browser) и руки (Telegram, Android ADB).\n• 7-этапная Конституционная Безопасность."

    # Слайд 4: Инновация - Self-Healing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Инновация: Meta-Cognitive Auto-Healing"
    content = slide.placeholders[1]
    content.text = "• Движок способен читать свой собственный код (AST Parsing).\n• При обнаружении критической ошибки (Crash), ИИ сам пишет патч.\n• Автономно коммитит исправления в GitHub.\n• Генерирует новые навыки (Zero-to-One) под запросы бизнеса."

    # Слайд 5: Что под капотом (Octopus Merge)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Архитектурное Слияние (Octopus Integration)"
    content = slide.placeholders[1]
    content.text = "• 240+ Встроенных Навыков.\n• P2P FastAPI децентрализация узлов.\n• Operator Matrix (WebSocket Дашборд + Next.js UI).\n• Мониторинг метрик роя в Grafana и Prometheus."

    # Слайд 6: Будущее (Внедрение)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AIOS: Готов к коммерческой эксплуатации"
    content = slide.placeholders[1]
    content.text = "• Commercial RPA Pipelines: автоматический парсинг лидов и отправка в Телеграм.\n• Готовый Docker Swarm стек (deploy_swarm.sh).\n\nБудущее наступило сегодня."

    prs.save("/home/user/AIOS/docs/AIOS_Pitch_Deck_v19.pptx")
    print("Pitch deck created at /home/user/AIOS/docs/AIOS_Pitch_Deck_v19.pptx")

if __name__ == "__main__":
    create_pitch_deck()
